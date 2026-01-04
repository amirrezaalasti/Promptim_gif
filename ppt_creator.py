"""PowerPoint presentation creation with embedded animations."""

import os
from pathlib import Path
from typing import Optional


def create_powerpoint_with_gif(
    gif_path: str,
    title: str = "Manim Animation",
    description: str = "",
    output_path: Optional[str] = None
) -> Optional[str]:
    """Create a PowerPoint presentation with embedded GIF.
    
    Note: PowerPoint doesn't natively animate GIFs in the editor,
    but they will animate during slideshow presentation.
    
    Args:
        gif_path: Path to the GIF file
        title: Slide title
        description: Description text
        output_path: Output path for PPTX (auto-generated if None)
        
    Returns:
        Path to generated PPTX or None on failure
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise ImportError("Please install python-pptx: pip install python-pptx")
    
    if output_path is None:
        output_path = gif_path.replace(".gif", ".pptx")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Title slide
    title_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(title_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(15), Inches(1.2)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title[:100]  # Limit title length
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(0x2E, 0x86, 0xAB)  # Nice blue
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add GIF (centered on slide)
    try:
        from PIL import Image
        
        img = Image.open(gif_path)
        img_width, img_height = img.size
        
        # Calculate dimensions to fit slide
        max_width = Inches(14)
        max_height = Inches(6.5)
        
        aspect = img_width / img_height
        
        if img_width / max_width.inches > img_height / max_height.inches:
            width = max_width
            height = Inches(max_width.inches / aspect)
        else:
            height = max_height
            width = Inches(max_height.inches * aspect)
        
        # Center position
        left = (prs.slide_width - width) / 2
        top = Inches(1.8)
        
        slide.shapes.add_picture(gif_path, left, top, width, height)
        
    except ImportError:
        # Fallback without PIL
        slide.shapes.add_picture(
            gif_path,
            Inches(1), Inches(1.8),
            Inches(14), Inches(6.5)
        )
    
    # Add description if provided
    if description:
        desc_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(8.2),
            Inches(15), Inches(0.6)
        )
        desc_frame = desc_box.text_frame
        desc_para = desc_frame.paragraphs[0]
        desc_para.text = description[:200]  # Limit description
        desc_para.font.size = Pt(14)
        desc_para.font.color.rgb = RgbColor(0x66, 0x66, 0x66)
        desc_para.alignment = PP_ALIGN.CENTER
    
    prs.save(output_path)
    return output_path


def create_multi_slide_presentation(
    animations: list,
    output_path: str,
    presentation_title: str = "Math Animations"
) -> Optional[str]:
    """Create a PowerPoint with multiple animation slides.
    
    Args:
        animations: List of dicts with 'gif_path', 'title', 'description'
        output_path: Output path for PPTX
        presentation_title: Title for the presentation
        
    Returns:
        Path to generated PPTX or None on failure
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise ImportError("Please install python-pptx: pip install python-pptx")
    
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Title slide
    title_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_layout)
    
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.5),
        Inches(14), Inches(2)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = presentation_title
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(0x2E, 0x86, 0xAB)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add animation slides
    for anim in animations:
        gif_path = anim.get("gif_path")
        title = anim.get("title", "Animation")
        description = anim.get("description", "")
        
        if not gif_path or not os.path.exists(gif_path):
            continue
        
        slide = prs.slides.add_slide(title_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(15), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title[:80]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RgbColor(0x2E, 0x86, 0xAB)
        title_para.alignment = PP_ALIGN.CENTER
        
        # GIF
        try:
            from PIL import Image
            
            img = Image.open(gif_path)
            img_width, img_height = img.size
            aspect = img_width / img_height
            
            max_width = Inches(14)
            max_height = Inches(6)
            
            if img_width / max_width.inches > img_height / max_height.inches:
                width = max_width
                height = Inches(max_width.inches / aspect)
            else:
                height = max_height
                width = Inches(max_height.inches * aspect)
            
            left = (prs.slide_width - width) / 2
            top = Inches(1.5)
            
            slide.shapes.add_picture(gif_path, left, top, width, height)
        except:
            slide.shapes.add_picture(
                gif_path,
                Inches(1), Inches(1.5),
                Inches(14), Inches(6)
            )
        
        # Description
        if description:
            desc_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(8),
                Inches(15), Inches(0.8)
            )
            desc_frame = desc_box.text_frame
            desc_para = desc_frame.paragraphs[0]
            desc_para.text = description[:200]
            desc_para.font.size = Pt(14)
            desc_para.font.color.rgb = RgbColor(0x66, 0x66, 0x66)
            desc_para.alignment = PP_ALIGN.CENTER
    
    prs.save(output_path)
    return output_path

